"""Build the Quantium Task 3 client presentation from the supplied brand template."""

from pathlib import Path
import shutil

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = Path.home() / "Downloads" / "Task 3 - presentation guide_BRAND.pptx"
TASK1_DATA = ROOT / "QVI_data.csv"
TRIAL_RESULTS = ROOT / "QVI_trial_results.csv"
OUT_PPTX = ROOT / "Quantium_Task3_Category_Review.pptx"
ASSET_DIR = ROOT / "task3_assets"
ASSET_DIR.mkdir(exist_ok=True)

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(0, 122, 128)
TEAL_DARK = RGBColor(0, 78, 84)
BLUE = RGBColor(47, 107, 138)
ORANGE = RGBColor(230, 126, 34)
RED = RGBColor(185, 55, 55)
GREY = RGBColor(90, 96, 102)
LIGHT = RGBColor(241, 244, 246)
GREEN = RGBColor(39, 125, 81)


def clear_slide(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def add_text(slide, text, x, y, w, h, size=18, color=BLACK, bold=False,
             align=PP_ALIGN.LEFT, font="Arial", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, text, subtitle=None):
    add_text(slide, text, .82, .33, 11.45, .58, 25, BLACK, True)
    if subtitle:
        add_text(slide, subtitle, .83, .92, 11.4, .35, 10.5, GREY)


def add_card(slide, x, y, w, h, headline, body, accent=TEAL, metric=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(220, 225, 228)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(.08), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    if metric:
        add_text(slide, metric, x+.22, y+.14, w-.38, .42, 20, accent, True)
        add_text(slide, headline, x+.22, y+.59, w-.38, .38, 12, BLACK, True)
        add_text(slide, body, x+.22, y+.98, w-.38, h-1.08, 9.5, GREY)
    else:
        add_text(slide, headline, x+.22, y+.17, w-.38, .42, 13, accent, True)
        add_text(slide, body, x+.22, y+.65, w-.38, h-.78, 10, GREY)


def add_footer(slide, page):
    # Page numbers are supplied by the branded slide master.
    return None


def add_picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def style_ax(ax, title=None):
    ax.spines[['top', 'right']].set_visible(False)
    ax.grid(axis='x', alpha=.18)
    ax.tick_params(labelsize=9)
    if title:
        ax.set_title(title, loc='left', fontsize=12, fontweight='bold')


def create_assets():
    d = pd.read_csv(TASK1_DATA, parse_dates=["DATE"])
    seg = d.groupby(["LIFESTAGE", "PREMIUM_CUSTOMER"]).agg(
        sales=("TOT_SALES", "sum"), customers=("LYLTY_CARD_NBR", "nunique"), units=("PROD_QTY", "sum")
    ).reset_index()
    seg["segment"] = seg.LIFESTAGE.str.title() + " | " + seg.PREMIUM_CUSTOMER
    seg["units_per_customer"] = seg.units / seg.customers
    seg["price_per_unit"] = seg.sales / seg.units

    # Slide 5: ranked sales pools.
    top = seg.nlargest(8, "sales").sort_values("sales")
    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    colors = ['#007A80' if 'Young Singles/Couples | Mainstream' in s else '#2F6B8A' for s in top.segment]
    ax.barh(top.segment, top.sales/1000, color=colors)
    for i, v in enumerate(top.sales/1000): ax.text(v+1.5, i, f"${v:.0f}k", va='center', fontsize=9)
    ax.set_xlabel("Annual chip sales ($000)"); style_ax(ax)
    fig.tight_layout(); p1=ASSET_DIR/'segment_sales.png'; fig.savefig(p1,dpi=180,transparent=True,bbox_inches='tight'); plt.close(fig)

    # Slide 6: sales drivers for the three largest pools.
    focus = seg.nlargest(3, "sales").copy()
    fig, axes = plt.subplots(1,2,figsize=(10,4.2))
    labels=[s.replace(' | ','\n') for s in focus.segment]
    axes[0].bar(labels, focus.units_per_customer, color=['#2F6B8A','#007A80','#6A7FDB'])
    axes[0].set_ylabel('Units per customer'); axes[0].tick_params(axis='x',rotation=0,labelsize=8); axes[0].set_title('Volume per buyer',loc='left',fontweight='bold')
    axes[1].bar(labels, focus.price_per_unit, color=['#2F6B8A','#007A80','#6A7FDB'])
    axes[1].set_ylabel('$ per unit'); axes[1].tick_params(axis='x',rotation=0,labelsize=8); axes[1].set_title('Average price paid',loc='left',fontweight='bold')
    for ax in axes: ax.spines[['top','right']].set_visible(False); ax.grid(axis='y',alpha=.18)
    fig.tight_layout(); p2=ASSET_DIR/'sales_drivers.png'; fig.savefig(p2,dpi=180,transparent=True,bbox_inches='tight'); plt.close(fig)

    # Slide 7: affinity of the target segment.
    target_mask = d.LIFESTAGE.eq('YOUNG SINGLES/COUPLES') & d.PREMIUM_CUSTOMER.eq('Mainstream')
    def affinity(field):
        t=d[target_mask].groupby(field).PROD_QTY.sum()/d[target_mask].PROD_QTY.sum()
        o=d[~target_mask].groupby(field).PROD_QTY.sum()/d[~target_mask].PROD_QTY.sum()
        return (t/o).dropna().sort_values(ascending=False)
    brands=affinity('BRAND').head(6).sort_values(); packs=affinity('PACK_SIZE').head(6).sort_values()
    fig, axes=plt.subplots(1,2,figsize=(10,4.2))
    axes[0].barh(brands.index,brands.values,color='#007A80'); axes[0].axvline(1,color='#777',lw=1); axes[0].set_xlabel('Affinity vs other shoppers'); axes[0].set_title('Preferred brands',loc='left',fontweight='bold')
    axes[1].barh([f'{x}g' for x in packs.index],packs.values,color='#E67E22'); axes[1].axvline(1,color='#777',lw=1); axes[1].set_xlabel('Affinity vs other shoppers'); axes[1].set_title('Preferred pack sizes',loc='left',fontweight='bold')
    for ax in axes: ax.spines[['top','right']].set_visible(False); ax.grid(axis='x',alpha=.18)
    fig.tight_layout(); p3=ASSET_DIR/'target_affinity.png'; fig.savefig(p3,dpi=180,transparent=True,bbox_inches='tight'); plt.close(fig)

    # Slide 9: control-match scores.
    trial=pd.read_csv(TRIAL_RESULTS)
    fig,ax=plt.subplots(figsize=(8,3.6))
    labels=[f"Trial {int(r.TRIAL_STORE)} → Control {int(r.CONTROL_STORE)}" for r in trial.itertuples()]
    ax.barh(labels[::-1],trial.MATCH_SCORE.values[::-1],color='#2F6B8A')
    ax.set_xlim(0,1); ax.set_xlabel('Composite similarity score')
    for i,v in enumerate(trial.MATCH_SCORE.values[::-1]): ax.text(v+.015,i,f'{v:.2f}',va='center')
    style_ax(ax); fig.tight_layout(); p4=ASSET_DIR/'control_matches.png'; fig.savefig(p4,dpi=180,transparent=True,bbox_inches='tight'); plt.close(fig)

    # Slide 10: trial outcomes.
    fig,ax=plt.subplots(figsize=(8.5,4.3))
    x=np.arange(len(trial)); width=.25
    ax.bar(x-width,trial.SALES_UPLIFT*100,width,label='Sales',color='#2F6B8A')
    ax.bar(x,trial.CUSTOMER_UPLIFT*100,width,label='Customers',color='#007A80')
    ax.bar(x+width,trial.TXN_PER_CUSTOMER_UPLIFT*100,width,label='Transactions/customer',color='#E67E22')
    ax.axhline(0,color='#555',lw=.8); ax.set_xticks(x,[f"Store {int(s)}" for s in trial.TRIAL_STORE]); ax.set_ylabel('Trial-period uplift vs scaled control (%)'); ax.legend(frameon=False,ncol=3,loc='upper center'); ax.spines[['top','right']].set_visible(False); ax.grid(axis='y',alpha=.18)
    fig.tight_layout(); p5=ASSET_DIR/'trial_uplift.png'; fig.savefig(p5,dpi=180,transparent=True,bbox_inches='tight'); plt.close(fig)
    return seg, trial, [p1,p2,p3,p4,p5]


def build():
    seg, trial, assets = create_assets()
    prs=Presentation(TEMPLATE)
    # Cover.
    s=prs.slides[0]
    for sh in s.shapes:
        if not sh.has_text_frame: continue
        if sh.name.startswith('Title'): sh.text='Category review: Chips'
        elif sh.name.startswith('Subtitle'): sh.text='Customer strategy and trial recommendations'
        elif 'Text Placeholder' in sh.name: sh.text='August 2026'

    # Executive summary.
    s=prs.slides[2]; clear_slide(s); add_title(s,'Grow chips through targeted activation and selective trial rollout','Answer first: three actions for the next six months')
    add_card(s,.65,1.5,3.7,4.55,'Win the growth segment','Target Mainstream young singles/couples with brand-led secondary displays. Prioritise Tyrrells and Twisties, with added visibility in the lead-up to Christmas.',TEAL,'01')
    add_card(s,4.58,1.5,3.7,4.55,'Protect the volume base','Use value bundles for older and young families. Avoid blanket discounting: the target segment already pays a higher unit price.',BLUE,'02')
    add_card(s,8.5,1.5,3.7,4.55,'Scale proven layouts','Roll out the trial design in stores resembling 77 and 88. Investigate store 86 execution and pricing before another controlled test.',ORANGE,'03')
    add_footer(s,3)

    # Divider 1.
    s=prs.slides[3]
    for sh in s.shapes:
        if sh.has_text_frame and sh.text.strip()=='01': sh.text='01'
        elif sh.has_text_frame: sh.text='Category growth strategy'

    # Category overview.
    s=prs.slides[4]; clear_slide(s); add_title(s,'Three customer pools account for one-quarter of category sales','Budget older families lead sales; Mainstream young singles/couples are the clearest growth target')
    add_picture(s,assets[0],.95,1.35,7.15,4.95)
    add_card(s,8.45,1.45,3.55,1.3,'Largest sales pool','Budget older families: $157k, or 8.7% of annual chip sales.',BLUE)
    add_card(s,8.45,3.02,3.55,1.3,'Growth target','Mainstream young singles/couples: $148k and the largest buyer base.',TEAL)
    add_card(s,8.45,4.59,3.55,1.3,'Seasonal opportunity','Transactions rise before Christmas; use gondola ends or promotional displays ahead of the holiday closure.',ORANGE)
    add_footer(s,5)

    # Drivers.
    s=prs.slides[5]; clear_slide(s); add_title(s,'Family volume and young-mainstream pricing drive category spend','Sales = customers × units per customer × average price per unit')
    add_picture(s,assets[1],.85,1.4,7.87,4.7)
    add_card(s,8.95,1.5,3.05,1.55,'Family volume','Older and young families buy more units per customer—protect them with value and bundle mechanics.',BLUE)
    add_card(s,8.95,3.35,3.05,1.55,'Price resilience','Mainstream young and mid-age singles/couples pay significantly more per unit (p < 0.001).',TEAL)
    add_text(s,'Commercial implication: target offers by mission; do not use a category-wide price cut.',8.98,5.3,3.0,.7,11,ORANGE,True)
    add_footer(s,6)

    # Target affinity.
    s=prs.slides[6]; clear_slide(s); add_title(s,'Tyrrells and Twisties are the strongest target-segment plays','Affinity compares the target segment’s unit mix with all other chip buyers')
    add_picture(s,assets[2],.95,1.35,7.75,4.75)
    add_card(s,8.95,1.45,3.0,1.55,'Brand activation','Off-locate Tyrrells and Twisties in high-traffic discretionary space; keep price architecture premium.',TEAL)
    add_card(s,8.95,3.35,3.0,1.55,'Pack-size caveat','270g leads pack affinity, but only Twisties sells this size—do not generalise the effect to other brands.',ORANGE)
    add_text(s,'Measure incremental units, margin and repeat purchase against matched controls.',8.98,5.3,2.95,.65,11,TEAL_DARK,True)
    add_footer(s,7)

    # Divider 2.
    s=prs.slides[7]
    for sh in s.shapes:
        if sh.has_text_frame and sh.text.strip()=='02': sh.text='02'
        elif sh.has_text_frame: sh.text='Trial store performance'

    # Controls.
    s=prs.slides[8]; clear_slide(s); add_title(s,'Matched controls isolate trial impact from normal variation','Controls were selected on seven pre-trial months of sales and customer trends and scale')
    add_picture(s,assets[3],.95,1.55,6.5,3.65)
    add_card(s,7.75,1.45,4.2,1.2,'77 → 233','Very strong match across sales and customer behaviour.',TEAL)
    add_card(s,7.75,2.9,4.2,1.2,'86 → 155','Strong match; suitable for diagnosing customer-versus-sales response.',BLUE)
    add_card(s,7.75,4.35,4.2,1.2,'88 → 237','Good match; remaining uncertainty is reflected in the rollout safeguards.',ORANGE)
    add_text(s,'Method: equal-weight composite of Pearson correlation and month-normalised magnitude similarity.',.85,5.65,10.9,.45,10,GREY)
    add_footer(s,9)

    # Trial outcomes.
    s=prs.slides[9]; clear_slide(s); add_title(s,'Roll out in stores like 77 and 88; investigate store 86 before retesting','A successful store cleared statistical and commercial thresholds in at least two of three trial months')
    add_picture(s,assets[4],.85,1.45,7.15,4.7)
    add_card(s,8.28,1.38,3.78,1.28,'Store 77 — ROLL OUT','Sales +26%; customers +23%. Growth came from reach, not more transactions per shopper.',GREEN)
    add_card(s,8.28,2.95,3.78,1.28,'Store 86 — INVESTIGATE','Customers +14% but sales were not consistently material. Check pricing, deals and execution.',RED)
    add_card(s,8.28,4.52,3.78,1.28,'Store 88 — ROLL OUT','Sales +12%; customers +6%; transactions/customer +6%. Scale to similar stores with monitoring.',GREEN)
    add_footer(s,10)

    prs.save(OUT_PPTX)
    print(f'Wrote {OUT_PPTX}')


if __name__=='__main__':
    build()
